#!/usr/bin/env python3
"""
PPT自動生成ツール（レイアウト画像解析 + Markdown）
使い方: python3 make_pptx_from_md.py <レイアウト画像> <Markdownファイル> [出力.pptx]

1. レイアウト画像をClaude APIで解析してデザイン設定を抽出
2. Markdownファイルをパースしてスライド構造に変換
3. python-pptxでPPTXファイルを生成
"""

import sys
import json
import base64
import re
from pathlib import Path

import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_WIDTH  = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# レイアウト解析結果のキャッシュファイル名（同じ画像を再利用する場合はAPI呼び出しをスキップ）
LAYOUT_CACHE_SUFFIX = "_layout_cache.json"


# ---------------------------------------------------------------------------
# ユーティリティ
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_color: str) -> RGBColor:
    """#RRGGBBをRGBColorに変換する"""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def add_text_box(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=None,
                 align=PP_ALIGN.LEFT, font_name="メイリオ"):
    if color is None:
        color = RGBColor(0x26, 0x26, 0x26)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


# ---------------------------------------------------------------------------
# レイアウト画像解析（Claude API Vision）
# ---------------------------------------------------------------------------

def analyze_layout(image_path: Path) -> dict:
    """Claude APIでレイアウト画像を解析してデザイン設定を返す"""
    cache_path = image_path.parent / (image_path.stem + LAYOUT_CACHE_SUFFIX)
    if cache_path.exists():
        print(f"キャッシュを使用: {cache_path}")
        with open(cache_path, encoding="utf-8") as f:
            return json.load(f)

    client = anthropic.Anthropic()

    with open(image_path, "rb") as f:
        img_b64 = base64.standard_b64encode(f.read()).decode("utf-8")

    ext = image_path.suffix.lower()
    media_type = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
    }.get(ext, "image/png")

    prompt = """このPowerPointスライドのレイアウト画像を解析し、デザイン設定をJSON形式のみで返してください。
余分な説明や```は不要です。JSONオブジェクトだけを返してください。

{
  "background_color": "#RRGGBB",
  "header_bg_color": "#RRGGBB",
  "header_text_color": "#RRGGBB",
  "body_text_color": "#RRGGBB",
  "accent_color": "#RRGGBB",
  "has_bottom_bar": true,
  "bottom_bar_color": "#RRGGBB",
  "header_height_ratio": 0.16,
  "title_font_size": 32,
  "body_font_size": 20,
  "font_name": "フォント名",
  "layout_style": "header_bar"
}

layout_style の選択肢:
- "header_bar"  : 上部に色帯ヘッダーがある標準スタイル
- "full_color"  : 背景全体が色付きのスタイル
- "minimal"     : シンプルで装飾が少ないスタイル

色が判断できない場合は近似値で推定してください。"""

    print("Claude APIでレイアウトを解析中...")
    resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {"type": "base64", "media_type": media_type, "data": img_b64},
                },
                {"type": "text", "text": prompt},
            ],
        }],
    )

    text = resp.content[0].text.strip()
    # コードブロックが含まれていれば除去
    m = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if m:
        text = m.group(1).strip()

    layout = json.loads(text)

    # キャッシュに保存
    with open(cache_path, "w", encoding="utf-8") as f:
        json.dump(layout, f, ensure_ascii=False, indent=2)
    print(f"解析結果をキャッシュに保存: {cache_path}")

    return layout


# ---------------------------------------------------------------------------
# Markdownパーサー
# ---------------------------------------------------------------------------

def md_body_to_text(md: str) -> str:
    """Markdownの本文をプレーンテキスト（箇条書き• 付き）に変換する"""
    lines = []
    for line in md.split("\n"):
        line = re.sub(r"^#{1,6}\s+", "", line)          # 見出し除去
        line = re.sub(r"\*\*(.+?)\*\*", r"\1", line)   # 太字除去
        line = re.sub(r"\*(.+?)\*", r"\1", line)        # イタリック除去
        line = re.sub(r"`(.+?)`", r"\1", line)          # インラインコード除去
        if re.match(r"^[-*+] ", line):
            line = "• " + line[2:]
        lines.append(line)
    return "\n".join(lines)


def parse_markdown(md_path: Path) -> dict:
    """Markdownファイルをプレゼンテーション構造（dict）に変換する。

    対応フォーマット:
      ---
      title: タイトル
      subtitle: サブタイトル
      author: 作成者
      date: 日付
      ---

      ## スライドタイトル

      本文テキスト
      - 箇条書き

      ## まとめ <!-- summary -->

      まとめ本文
    """
    with open(md_path, encoding="utf-8") as f:
        content = f.read()

    cfg = {"title": "", "subtitle": "", "author": "", "date": "", "slides": []}

    # フロントマター（---...---）を解析
    fm_match = re.match(r"^---\s*\n([\s\S]*?)\n---\s*\n?", content)
    if fm_match:
        for line in fm_match.group(1).split("\n"):
            if ":" in line:
                k, _, v = line.partition(":")
                k, v = k.strip(), v.strip()
                if k in cfg:
                    cfg[k] = v
        content = content[fm_match.end():]

    # フロントマターがない場合はH1をタイトルとして使用
    if not cfg["title"]:
        m = re.search(r"^# (.+)$", content, re.MULTILINE)
        if m:
            cfg["title"] = m.group(1).strip()
            content = content[:m.start()] + content[m.end():]

    # H2でスライドに分割
    parts = re.split(r"^## (.+)$", content, flags=re.MULTILINE)

    # parts[0] はH2より前（無視）
    i = 1
    while i < len(parts) - 1:
        raw_title = parts[i].strip()
        raw_body  = parts[i + 1].strip()

        # <!-- summary --> タグで種別を判定
        slide_type = "content"
        if "<!-- summary -->" in raw_title or "<!-- summary -->" in raw_body:
            slide_type = "summary"
        raw_title = re.sub(r"<!--.*?-->", "", raw_title).strip()
        raw_body  = re.sub(r"<!--.*?-->", "", raw_body).strip()

        cfg["slides"].append({
            "type": slide_type,
            "title": raw_title,
            "body": md_body_to_text(raw_body),
        })
        i += 2

    return cfg


# ---------------------------------------------------------------------------
# スライド生成
# ---------------------------------------------------------------------------

def _header_h(L: dict):
    return SLIDE_HEIGHT * L.get("header_height_ratio", 0.16)


def _bottom_bar(slide, L: dict):
    if L.get("has_bottom_bar", True):
        col = hex_to_rgb(L.get("bottom_bar_color") or L["header_bg_color"])
        add_rect(slide, 0, Inches(7.2), SLIDE_WIDTH, Inches(0.3), col)


def make_title_slide(prs, cfg: dict, L: dict):
    """表紙スライドを作成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg  = hex_to_rgb(L["background_color"])
    hbg = hex_to_rgb(L["header_bg_color"])
    htx = hex_to_rgb(L["header_text_color"])
    btx = hex_to_rgb(L["body_text_color"])
    fn  = L.get("font_name", "メイリオ")
    tfs = L.get("title_font_size", 36)
    bfs = L.get("body_font_size", 22)

    set_bg(slide, bg)
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.15), hbg)

    # タイトル帯
    add_rect(slide, Inches(1), Inches(1.8), Inches(11.33), Inches(2.5), hbg)
    add_text_box(slide, cfg.get("title", ""),
                 Inches(1.3), Inches(2.0), Inches(10.8), Inches(1.2),
                 font_size=tfs, bold=True, color=htx,
                 align=PP_ALIGN.CENTER, font_name=fn)
    add_text_box(slide, cfg.get("subtitle", ""),
                 Inches(1.3), Inches(3.1), Inches(10.8), Inches(0.8),
                 font_size=bfs, color=htx, align=PP_ALIGN.CENTER, font_name=fn)

    info = f"{cfg.get('date', '')}　{cfg.get('author', '')}".strip("　 ")
    add_text_box(slide, info,
                 Inches(1), Inches(5.5), Inches(11.33), Inches(0.5),
                 font_size=14, color=btx, align=PP_ALIGN.RIGHT, font_name=fn)

    _bottom_bar(slide, L)


def make_toc_slide(prs, cfg: dict, L: dict):
    """目次スライドを作成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg  = hex_to_rgb(L["background_color"])
    hbg = hex_to_rgb(L["header_bg_color"])
    htx = hex_to_rgb(L["header_text_color"])
    btx = hex_to_rgb(L["body_text_color"])
    fn  = L.get("font_name", "メイリオ")
    hh  = _header_h(L)

    set_bg(slide, bg)
    add_rect(slide, 0, 0, SLIDE_WIDTH, hh, hbg)
    add_text_box(slide, "目次",
                 Inches(0.5), Inches(0.15), Inches(12), hh - Inches(0.1),
                 font_size=L.get("title_font_size", 32), bold=True,
                 color=htx, font_name=fn)

    titles = [s["title"] for s in cfg.get("slides", [])]
    body = "\n".join(f"  {i+1}.  {t}" for i, t in enumerate(titles))
    add_text_box(slide, body,
                 Inches(1.0), hh + Inches(0.2),
                 Inches(11), SLIDE_HEIGHT - hh - Inches(0.6),
                 font_size=L.get("body_font_size", 20), color=btx, font_name=fn)

    _bottom_bar(slide, L)


def make_content_slide(prs, slide_cfg: dict, L: dict):
    """本文スライドを作成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg  = hex_to_rgb(L["background_color"])
    hbg = hex_to_rgb(L["header_bg_color"])
    htx = hex_to_rgb(L["header_text_color"])
    btx = hex_to_rgb(L["body_text_color"])
    fn  = L.get("font_name", "メイリオ")
    hh  = _header_h(L)

    set_bg(slide, bg)
    add_rect(slide, 0, 0, SLIDE_WIDTH, hh, hbg)
    add_text_box(slide, slide_cfg["title"],
                 Inches(0.5), Inches(0.15), Inches(12), hh - Inches(0.1),
                 font_size=L.get("title_font_size", 32), bold=True,
                 color=htx, font_name=fn)
    add_text_box(slide, slide_cfg["body"],
                 Inches(0.8), hh + Inches(0.2),
                 Inches(11.7), SLIDE_HEIGHT - hh - Inches(0.6),
                 font_size=L.get("body_font_size", 20), color=btx, font_name=fn)

    _bottom_bar(slide, L)


def make_summary_slide(prs, slide_cfg: dict, L: dict):
    """まとめスライドを作成する（ヘッダー色を背景に使って強調）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hbg = hex_to_rgb(L["header_bg_color"])
    htx = hex_to_rgb(L["header_text_color"])
    fn  = L.get("font_name", "メイリオ")

    set_bg(slide, hbg)
    add_text_box(slide, slide_cfg["title"],
                 Inches(0.5), Inches(0.3), Inches(12), Inches(0.9),
                 font_size=L.get("title_font_size", 36), bold=True,
                 color=htx, align=PP_ALIGN.CENTER, font_name=fn)
    add_rect(slide, Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.05), htx)
    add_text_box(slide, slide_cfg["body"],
                 Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5),
                 font_size=L.get("body_font_size", 22), color=htx, font_name=fn)


# ---------------------------------------------------------------------------
# メイン
# ---------------------------------------------------------------------------

def main():
    if len(sys.argv) < 3:
        print("使い方: python3 make_pptx_from_md.py <レイアウト画像> <Markdownファイル> [出力.pptx]")
        print()
        print("例:")
        print("  python3 make_pptx_from_md.py layout.png content.md")
        print("  python3 make_pptx_from_md.py layout.png content.md output/my_slide.pptx")
        sys.exit(1)

    image_path  = Path(sys.argv[1])
    md_path     = Path(sys.argv[2])
    output_path = Path(sys.argv[3]) if len(sys.argv) >= 4 else md_path.with_suffix(".pptx")

    if not image_path.exists():
        print(f"エラー: 画像ファイルが見つかりません: {image_path}")
        sys.exit(1)
    if not md_path.exists():
        print(f"エラー: Markdownファイルが見つかりません: {md_path}")
        sys.exit(1)

    # ステップ1: レイアウト解析
    layout = analyze_layout(image_path)
    print(f"レイアウト設定: {json.dumps(layout, ensure_ascii=False)}")

    # ステップ2: Markdown解析
    print("Markdownを解析中...")
    cfg = parse_markdown(md_path)
    print(f"タイトル: {cfg['title']}")
    print(f"スライド数: {len(cfg['slides'])}枚")

    # ステップ3: PPTX生成
    print("PPTXを生成中...")
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    make_title_slide(prs, cfg, layout)

    if len(cfg.get("slides", [])) >= 2:
        make_toc_slide(prs, cfg, layout)

    for slide_cfg in cfg.get("slides", []):
        if slide_cfg.get("type") == "summary":
            make_summary_slide(prs, slide_cfg, layout)
        else:
            make_content_slide(prs, slide_cfg, layout)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"\n生成完了: {output_path}")
    print(f"総スライド数: {prs.slides.__len__()}枚（表紙 + 目次 + 本文）")


if __name__ == "__main__":
    main()
