#!/usr/bin/env python3
"""
tenkizu README.md → PowerPoint 生成スクリプト
テンプレートPPTX（outline_info2026-2.pptx）の配色を使用
スライドサイズ: 4:3（10×7.5インチ、テンプレートに合わせる）
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================
# テンプレートから抽出した配色（theme1.xml準拠）
# ============================================================
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)   # 背景: 白
C_HEADER   = RGBColor(0x33, 0x33, 0x99)   # ヘッダー帯: ネイビー（accent2）
C_ACCENT   = RGBColor(0xBB, 0xE0, 0xE3)   # アクセント: ライトブルー（accent1）
C_ACCENT2  = RGBColor(0x2D, 0x2D, 0x8A)   # サブアクセント: 濃紺（accent6）
C_TEAL     = RGBColor(0x00, 0x99, 0x99)   # ティール（hlink）
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)   # 白
C_BLACK    = RGBColor(0x00, 0x00, 0x00)   # 黒（本文）
C_GRAY     = RGBColor(0x40, 0x40, 0x40)   # ダークグレー（サブ本文）

# スライドサイズ（4:3）
W = Inches(10)
H = Inches(7.5)

FONT = "メイリオ"


# ============================================================
# ユーティリティ
# ============================================================

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, line=False):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
    else:
        shp.line.fill.background()
    return shp


def txbox(slide, text, l, t, w, h,
          fs=18, bold=False, color=C_BLACK,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return box


def header_bar(slide, title, fs=22):
    """上部ヘッダー帯 + タイトルテキスト"""
    hh = Inches(1.05)
    rect(slide, 0, 0, W, hh, C_HEADER)
    txbox(slide, title, Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.85),
          fs=fs, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)


def slide_num_label(slide, num):
    """右下スライド番号"""
    txbox(slide, str(num), W - Inches(0.5), H - Inches(0.35), Inches(0.4), Inches(0.3),
          fs=11, color=C_GRAY, align=PP_ALIGN.RIGHT)


def content_start_y():
    return Inches(1.2)


def bottom_line(slide):
    """下部アクセントライン"""
    rect(slide, 0, H - Inches(0.12), W, Inches(0.12), C_ACCENT)


# ============================================================
# 個別スライド生成
# ============================================================

def slide_01_title(prs):
    """表紙"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)

    # 上部帯
    rect(sld, 0, 0, W, Inches(0.2), C_HEADER)

    # 中央タイトル帯
    rect(sld, Inches(0.5), Inches(1.8), W - Inches(1), Inches(2.6), C_HEADER)

    txbox(sld, "tenkizu",
          Inches(0.8), Inches(2.0), W - Inches(1.6), Inches(0.85),
          fs=42, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    txbox(sld, "天気図作成ツール（GSM / ECMWF）",
          Inches(0.8), Inches(2.75), W - Inches(1.6), Inches(0.7),
          fs=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    txbox(sld, "GSM・ECMWF GRIB2データから各種高層・地上天気図を自動生成するPythonツール",
          Inches(0.8), Inches(3.4), W - Inches(1.6), Inches(0.7),
          fs=16, color=C_WHITE, align=PP_ALIGN.CENTER)

    txbox(sld, "Python 3.10  |  conda: met_env_310",
          Inches(0.5), Inches(5.5), W - Inches(1), Inches(0.5),
          fs=14, color=C_GRAY, align=PP_ALIGN.RIGHT)

    rect(sld, 0, H - Inches(0.25), W, Inches(0.25), C_ACCENT2)


def slide_02_overview(prs, num):
    """概要"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "概要")
    bottom_line(sld)

    items = [
        ("対応データ",   "気象庁 GSM（全球モデル）・ECMWF（欧州中期予報センター）"),
        ("スクリプト数", "描画スクリプト 16種類（GSM 11本・ECM 6本）＋レポート生成 3本"),
        ("図法",         "ステレオ投影（中心: 60°N, 140°E）"),
        ("描画領域",     "108〜156°E, 17〜55°N（極東域）"),
        ("出力形式",     "PNG（DPI 150、10×8インチ）"),
        ("実行環境",     "Python 3.10（conda 環境 met_env_310）"),
    ]

    y = content_start_y()
    row_h = Inches(0.85)
    for i, (label, desc) in enumerate(items):
        top = y + i * row_h
        rect(sld, Inches(0.4), top, Inches(2.2), row_h - Inches(0.08), C_HEADER)
        txbox(sld, label, Inches(0.45), top + Inches(0.18), Inches(2.1), Inches(0.55),
              fs=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txbox(sld, desc, Inches(2.75), top + Inches(0.15), Inches(6.9), Inches(0.6),
              fs=14, color=C_BLACK)

    slide_num_label(sld, num)


def slide_03_gsm(prs, num):
    """データソース: GSM"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "データソース ①: GSM（気象庁 全球モデル）")
    bottom_line(sld)

    rows = [
        ("提供元",         "京都大学生存圏研究所 (RISH) データベース"),
        ("更新頻度",       "1日2回（00UTC・12UTC）"),
        ("利用可能期間",   "過去データも無償取得可（長期アーカイブ）"),
        ("ファイル形式",   "GRIB2（.bin 拡張子）"),
        ("水平解像度",     "約13km（0.125°）"),
        ("予報時間",       "0〜264h（〜72h: 6h間隔、72h〜: 12h間隔）"),
        ("収録変数",       "gh, t, u, v, w, r（高層）/ prmsl, sp, 10m風, 2m気温, 雲量"),
        ("注意点",         "可降水量(tcwv)・積算降水量(tp) は Rglファイルに含まれない"),
    ]

    y = content_start_y()
    row_h = Inches(0.72)
    for i, (label, desc) in enumerate(rows):
        top = y + i * row_h
        bg = C_ACCENT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), top, W - Inches(0.6), row_h - Inches(0.05), bg)
        txbox(sld, label, Inches(0.35), top + Inches(0.1), Inches(1.8), Inches(0.5),
              fs=13, bold=True, color=C_HEADER)
        txbox(sld, desc, Inches(2.25), top + Inches(0.1), Inches(7.4), Inches(0.5),
              fs=13, color=C_BLACK)

    slide_num_label(sld, num)


def slide_04_ecmwf(prs, num):
    """データソース: ECMWF"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "データソース ②: ECMWF（欧州中期予報センター）")
    bottom_line(sld)

    rows = [
        ("提供元",         "ECMWF Open Data（無償公開）"),
        ("更新頻度",       "1日4回（00/06/12/18UTC）"),
        ("利用可能期間",   "最新約5日分のみ無償。過去データは Copernicus CDS API が必要"),
        ("ファイル形式",   "GRIB2（.grib2 拡張子）"),
        ("水平解像度",     "約9km（0.25°）"),
        ("予報時間",       "0〜240h（〜144h: 3h間隔、144h〜: 6h間隔）"),
        ("収録変数",       "高層変数に加え tcwv（可降水量）・tp（積算降水量）・skt 等も収録"),
        ("ライセンス",     "CC BY 4.0"),
    ]

    y = content_start_y()
    row_h = Inches(0.72)
    for i, (label, desc) in enumerate(rows):
        top = y + i * row_h
        bg = C_ACCENT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), top, W - Inches(0.6), row_h - Inches(0.05), bg)
        txbox(sld, label, Inches(0.35), top + Inches(0.1), Inches(1.8), Inches(0.5),
              fs=13, bold=True, color=C_HEADER)
        txbox(sld, desc, Inches(2.25), top + Inches(0.1), Inches(7.4), Inches(0.5),
              fs=13, color=C_BLACK)

    slide_num_label(sld, num)


def slide_05_comparison(prs, num):
    """GSM vs ECMWF 比較"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "GSM vs ECMWF — 比較と使い分け")
    bottom_line(sld)

    y0 = content_start_y()

    # ヘッダー行
    cols = [Inches(0.3), Inches(2.8), Inches(5.8), Inches(8.0)]
    col_w = [Inches(2.45), Inches(2.95), Inches(2.15), Inches(1.85)]
    headers = ["項目", "GSM（気象庁）", "ECMWF", "優位"]

    row_h = Inches(0.55)

    for j, (hdr, cw, cx) in enumerate(zip(headers, col_w, cols)):
        rect(sld, cx, y0, cw - Inches(0.05), row_h, C_HEADER)
        txbox(sld, hdr, cx + Inches(0.05), y0 + Inches(0.08), cw - Inches(0.1), Inches(0.4),
              fs=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("水平解像度",  "約13km",          "約9km",           "ECM"),
        ("予報時間",    "〜264h",          "〜240h",          "GSM"),
        ("更新頻度",    "2回/日",          "4回/日",          "ECM"),
        ("無償取得",    "過去データも全て", "最新5日分のみ",   "GSM"),
        ("地表面変数",  "限定的",          "豊富(tcwv/tp等)",  "ECM"),
        ("ライセンス",  "気象庁利用規約",  "CC BY 4.0",       "—"),
    ]

    for i, (item, gsm, ecm, win) in enumerate(rows):
        top = y0 + row_h + i * row_h
        bg = RGBColor(0xF5, 0xF5, 0xFF) if i % 2 == 0 else C_WHITE
        for cx, cw in zip(cols, col_w):
            rect(sld, cx, top, cw - Inches(0.05), row_h - Inches(0.03), bg)

        vals = [item, gsm, ecm, win]
        for j, (val, cw, cx) in enumerate(zip(vals, col_w, cols)):
            col_c = C_HEADER if j == 0 else (C_TEAL if val == win else C_BLACK)
            txbox(sld, val, cx + Inches(0.05), top + Inches(0.08), cw - Inches(0.1), Inches(0.4),
                  fs=12, bold=(j == 0), color=col_c, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    # 用途メモ
    y_note = y0 + row_h * 7 + Inches(0.1)
    txbox(sld, "● 過去事例解析 → GSM（RISHアーカイブで長期データ）　"
               "● 最新高精度予報 → ECMWF　● 可降水量・積算降水量 → ECMWF のみ",
          Inches(0.3), y_note, W - Inches(0.6), Inches(0.45),
          fs=12, italic=True, color=C_ACCENT2)

    slide_num_label(sld, num)


def slide_06_scripts_gsm(prs, num):
    """GSM系スクリプト一覧"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "GSM系スクリプト一覧（11本）")
    bottom_line(sld)

    scripts = [
        ("GSM_tenkizu500hPa.py",  "500hPa 等高度線・渦度シェード・H/L"),
        ("GSM_QVector850hPa.py",  "850hPa Qベクター発散・等温度線・等高度線"),
        ("GSM_Jet300hPa.py",      "300hPa 等風速線・非地衡風・収束発散"),
        ("GSM_Instability.py",    "大気不安定域分布（SEPT−maxEPT差）"),
        ("GSM_CrossSection.py",   "ポテンシャル温位・EPT・風の鉛直断面"),
        ("GSM_fax57.py",          "FAX57相当（500hPa気温・700hPa湿数）"),
        ("GSM_fax78.py",          "FAX78相当（850hPa気温・風・700hPa発散）"),
        ("GSM_faxSrfPre.py",      "地上気圧（等圧線）・10m風矢羽・2m等温度線"),
        ("GSM_EPT850hPa.py",      "850hPa 相当温位シェード・等値線・風矢羽"),
        ("GSM_100hPa.py",         "任意気圧面 等高度線・ISOTAC・風矢羽（100hPa等）"),
        ("GSM_tenkizu500hPa",     "旧メイン版後継（kurora_tenkizu.py 互換）"),
    ]

    y = content_start_y()
    row_h = Inches(0.53)
    for i, (name, desc) in enumerate(scripts):
        top = y + i * row_h
        bg = C_ACCENT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), top, W - Inches(0.6), row_h - Inches(0.04), bg)
        txbox(sld, name, Inches(0.35), top + Inches(0.07), Inches(3.0), Inches(0.42),
              fs=11, bold=True, color=C_HEADER, italic=True)
        txbox(sld, desc, Inches(3.45), top + Inches(0.07), Inches(6.15), Inches(0.42),
              fs=12, color=C_BLACK)

    slide_num_label(sld, num)


def slide_07_scripts_ecm(prs, num):
    """ECMWF系スクリプト一覧"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "ECMWF系スクリプト一覧（6本）")
    bottom_line(sld)

    scripts = [
        ("ECM_tenkizu500hPa.py",   "500hPa 等高度線・渦度シェード・H/L"),
        ("ECM_EPT850hPa.py",       "850hPa 相当温位シェード・等値線・風矢羽"),
        ("ECM_Fax57.py",           "FAX57相当（500hPa気温・700hPa湿数）"),
        ("ECM_Fax78.py",           "FAX78相当（850hPa気温・風・700hPa発散）"),
        ("ECM_SurfacePressure.py", "地上気圧・10m風矢羽・2m等温度線（±tcwv/tp オプション）"),
        ("ECM_100hPa.py",          "任意気圧面 等高度線・ISOTAC・風矢羽（100hPa等）"),
    ]

    y = content_start_y()
    row_h = Inches(0.7)
    for i, (name, desc) in enumerate(scripts):
        top = y + i * row_h
        bg = C_ACCENT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), top, W - Inches(0.6), row_h - Inches(0.06), bg)
        txbox(sld, name, Inches(0.35), top + Inches(0.1), Inches(3.0), Inches(0.5),
              fs=12, bold=True, color=C_HEADER, italic=True)
        txbox(sld, desc, Inches(3.45), top + Inches(0.1), Inches(6.15), Inches(0.5),
              fs=13, color=C_BLACK)

    # ECMWFの注意事項
    note_y = y + row_h * 6 + Inches(0.15)
    rect(sld, Inches(0.3), note_y, W - Inches(0.6), Inches(0.65), RGBColor(0xFF, 0xF0, 0xD0))
    txbox(sld, "⚠ ECMWF Open Data は最新約5日分のみ無償。過去データは Copernicus CDS API が必要。\n"
               "   START_FT は時間数（0, 6, 24 …）で指定（GSMのDDHH形式ではない）",
          Inches(0.4), note_y + Inches(0.05), W - Inches(0.8), Inches(0.55),
          fs=12, color=C_ACCENT2)

    slide_num_label(sld, num)


def slide_08_args_and_bulk(prs, num):
    """引数仕様と一括実行"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "引数仕様と一括実行")
    bottom_line(sld)

    y = content_start_y()

    # 共通引数
    txbox(sld, "■ 共通引数仕様",
          Inches(0.3), y, W - Inches(0.6), Inches(0.35), fs=14, bold=True, color=C_HEADER)
    args = [
        ("INIT_TIME", "必須", "初期時刻 YYYYMMDDHH（UTC）"),
        ("START_FT",  "省略可", "開始予報時間。GSM: DDHH形式  /  ECM: 時間数"),
        ("N_STEPS",   "省略可", "作成枚数（6h間隔）。省略時は1枚"),
    ]
    for i, (arg, req, desc) in enumerate(args):
        top = y + Inches(0.4) + i * Inches(0.5)
        rect(sld, Inches(0.3), top, W - Inches(0.6), Inches(0.46), C_ACCENT if i % 2 == 0 else C_WHITE)
        txbox(sld, arg, Inches(0.4), top + Inches(0.08), Inches(1.5), Inches(0.35),
              fs=12, bold=True, color=C_HEADER, italic=True)
        txbox(sld, req, Inches(1.95), top + Inches(0.08), Inches(0.9), Inches(0.35),
              fs=12, color=C_TEAL, align=PP_ALIGN.CENTER)
        txbox(sld, desc, Inches(2.9), top + Inches(0.08), Inches(6.7), Inches(0.35), fs=12, color=C_BLACK)

    # 一括実行
    sep_y = y + Inches(2.0)
    txbox(sld, "■ 一括実行: run_all_charts.sh",
          Inches(0.3), sep_y, W - Inches(0.6), Inches(0.35), fs=14, bold=True, color=C_HEADER)

    cmds = [
        "bash run_all_charts.sh 2026041200                # FT=0h 各1枚",
        "bash run_all_charts.sh 2026041200 0000 5         # FT=0〜24h 各5枚",
        "bash run_all_charts.sh 2026041200 0000 key       # FT=0/12/24/36/48h（keyモード）",
        "bash run_all_charts.sh 2026041200 --ecm          # ECM系も追加実行",
    ]
    code_y = sep_y + Inches(0.4)
    rect(sld, Inches(0.3), code_y, W - Inches(0.6), Inches(1.7), RGBColor(0xF0, 0xF0, 0xF8))
    for i, cmd in enumerate(cmds):
        txbox(sld, cmd, Inches(0.45), code_y + Inches(0.1) + i * Inches(0.38),
              W - Inches(0.9), Inches(0.35), fs=11, color=C_ACCENT2, italic=True)

    # keyモード注記
    key_y = code_y + Inches(1.75)
    txbox(sld, "keyモード: FT=0/12/24/36/48h の5枚を自動生成（ECM系引数変換も内部で自動処理）",
          Inches(0.3), key_y, W - Inches(0.6), Inches(0.35), fs=12, italic=True, color=C_GRAY)

    slide_num_label(sld, num)


def slide_09_auto(prs, num):
    """自動データ取得・一括生成"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "自動データ取得 ＆ 一括生成スクリプト")
    bottom_line(sld)

    y = content_start_y()

    def script_block(title, color, cmds, options, top):
        rect(sld, Inches(0.3), top, W - Inches(0.6), Inches(0.38), color)
        txbox(sld, title, Inches(0.4), top + Inches(0.06), W - Inches(0.8), Inches(0.3),
              fs=14, bold=True, color=C_WHITE)
        cmd_y = top + Inches(0.42)
        rect(sld, Inches(0.3), cmd_y, W - Inches(0.6), Inches(len(cmds) * 0.36 + 0.1),
             RGBColor(0xF0, 0xF0, 0xF8))
        for i, cmd in enumerate(cmds):
            txbox(sld, cmd, Inches(0.45), cmd_y + Inches(0.05) + i * Inches(0.36),
                  W - Inches(0.9), Inches(0.32), fs=10, color=C_ACCENT2, italic=True)
        opt_y = cmd_y + Inches(len(cmds) * 0.36 + 0.15)
        for i, (opt, desc) in enumerate(options):
            txbox(sld, f"  {opt}", Inches(0.3), opt_y + i * Inches(0.36),
                  Inches(2.3), Inches(0.32), fs=11, bold=True, color=C_HEADER, italic=True)
            txbox(sld, desc, Inches(2.65), opt_y + i * Inches(0.36),
                  Inches(7.0), Inches(0.32), fs=11, color=C_BLACK)
        return opt_y + Inches(len(options) * 0.36 + 0.2)

    gsm_cmds = [
        "python run_gsm_auto.py                              # 最新データ・keyモード",
        "python run_gsm_auto.py --steps 5                   # 最新データ・5枚連続",
        "python run_gsm_auto.py --init-time 2026041200      # 初期時刻を手動指定",
    ]
    gsm_opts = [
        ("--init-time", "初期時刻 YYYYMMDDHH（省略時はRISHサーバーから自動検索）"),
        ("--steps",     "連続枚数（6h間隔。省略時はkeyモード）"),
        ("--start-ft",  "開始予報時間（DDHH形式、--steps使用時）"),
    ]
    next_y = script_block("run_gsm_auto.py — GSM系 最新データ自動検索・一括実行",
                          C_HEADER, gsm_cmds, gsm_opts, y)

    ecm_cmds = [
        "python run_ecm_auto.py                             # 最新データ・keyモード",
        "python run_ecm_auto.py --tcwv                      # 地上気圧図に可降水量追加",
        "python run_ecm_auto.py --tp                        # 地上気圧図に積算降水量追加（FT>0必須）",
    ]
    ecm_opts = [
        ("--init-time", "初期時刻（省略時はECMWFサーバーHEADリクエストで自動検索）"),
        ("--tcwv/--tp", "ECM_SurfacePressure.py へのオプション転送"),
    ]
    script_block("run_ecm_auto.py — ECMWF系 最新データ自動検索・一括実行",
                 C_ACCENT2, ecm_cmds, ecm_opts, next_y)

    slide_num_label(sld, num)


def slide_10_reports(prs, num):
    """レポート生成スクリプト"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "レポート生成スクリプト")
    bottom_line(sld)

    y = content_start_y()

    scripts = [
        ("upper_wind_report.py",
         "上層天気図レポート生成",
         "指定気圧面（デフォルト100hPa）の等高度線・ISOTAC・風矢羽を生成し、"
         "reports/{init_str}/ にPNGとMarkdownをまとめて保存。--push でGitHub push。",
         "python upper_wind_report.py 2026041200 0000 5 --levels 100 50 --ecm --push"),
        ("synop_report.py",
         "総観天気図レポート生成",
         "Jet300hPa・Fax57・Fax78・850hPa EPT・地上気圧を組み合わせた"
         "チャートセットを一括生成し、Markdownレポートにまとめる。",
         "python synop_report.py 2026041200 0000 5 --ecm --push"),
        ("jet_front_report.py",
         "ジェット・前線解析レポート",
         "上層風・断面図・850hPa EPT・地上気圧の解析レポートを生成。",
         "python jet_front_report.py 2026041200 0000 5 --push"),
    ]

    row_h = Inches(1.7)
    for i, (script, title, desc, cmd) in enumerate(scripts):
        top = y + i * row_h
        rect(sld, Inches(0.3), top, W - Inches(0.6), row_h - Inches(0.1),
             RGBColor(0xF5, 0xF5, 0xFF))
        txbox(sld, f"■ {script}", Inches(0.4), top + Inches(0.05), W - Inches(0.8), Inches(0.3),
              fs=13, bold=True, color=C_HEADER, italic=True)
        txbox(sld, title, Inches(0.4), top + Inches(0.35), Inches(3.5), Inches(0.28),
              fs=12, bold=True, color=C_ACCENT2)
        txbox(sld, desc, Inches(0.4), top + Inches(0.65), W - Inches(0.8), Inches(0.45),
              fs=11, color=C_BLACK)
        rect(sld, Inches(0.4), top + row_h - Inches(0.55), W - Inches(0.8), Inches(0.38),
             RGBColor(0xE8, 0xE8, 0xF8))
        txbox(sld, cmd, Inches(0.5), top + row_h - Inches(0.5), W - Inches(1.0), Inches(0.3),
              fs=10, color=C_ACCENT2, italic=True)

    slide_num_label(sld, num)


def slide_11_pptx_gen(prs, num):
    """PowerPoint自動生成"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "PowerPoint 自動生成（make_pptx.py / make_pptx2.py）")
    bottom_line(sld)

    y = content_start_y()

    txbox(sld, "output/ 内のPNG画像をPowerPointに自動貼り付け。アスペクト比を保ったままセル内に収める（Pillow使用）",
          Inches(0.3), y, W - Inches(0.6), Inches(0.4), fs=13, color=C_GRAY)

    # make_pptx.py
    y1 = y + Inches(0.5)
    rect(sld, Inches(0.3), y1, W - Inches(0.6), Inches(0.35), C_HEADER)
    txbox(sld, "make_pptx.py — 主要7グループ",
          Inches(0.4), y1 + Inches(0.05), W - Inches(0.8), Inches(0.28),
          fs=13, bold=True, color=C_WHITE)

    groups1 = [
        ("GSM: 500hPa渦度/地上気圧", "2×2"),
        ("GSM: Fax57/Fax78",         "2×2"),
        ("GSM: 300hPaジェット/850hPa Qベクター", "2×2"),
        ("GSM: 850hPa相当温位",       "4-in-1"),
        ("ECM: 500hPa渦度/地上気圧",  "2×2"),
        ("ECM: Fax57/Fax78",          "2×2"),
        ("ECM: 850hPa相当温位",        "4-in-1"),
    ]
    gr_y = y1 + Inches(0.4)
    gr_h = Inches(0.35)
    for i, (name, mode) in enumerate(groups1):
        top = gr_y + i * gr_h
        bg = C_ACCENT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), top, Inches(7.5), gr_h - Inches(0.03), bg)
        txbox(sld, name, Inches(0.4), top + Inches(0.05), Inches(6.5), Inches(0.27), fs=11, color=C_BLACK)
        rect(sld, Inches(7.85), top, Inches(1.8), gr_h - Inches(0.03), C_ACCENT2)
        txbox(sld, mode, Inches(7.85), top + Inches(0.05), Inches(1.8), Inches(0.27),
              fs=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # make_pptx2.py
    y2 = gr_y + len(groups1) * gr_h + Inches(0.25)
    rect(sld, Inches(0.3), y2, W - Inches(0.6), Inches(0.35), C_TEAL)
    txbox(sld, "make_pptx2.py — 補完3グループ",
          Inches(0.4), y2 + Inches(0.05), W - Inches(0.8), Inches(0.28),
          fs=13, bold=True, color=C_WHITE)

    groups2 = [
        ("GSM: 300hPaジェット/Fax57", "2×2"),
        ("GSM: 大気不安定域/850hPa相当温位", "2×2"),
        ("GSM: 鉛直断面図",            "1×2"),
    ]
    gr2_y = y2 + Inches(0.4)
    for i, (name, mode) in enumerate(groups2):
        top = gr2_y + i * gr_h
        bg = RGBColor(0xD0, 0xF0, 0xF0) if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), top, Inches(7.5), gr_h - Inches(0.03), bg)
        txbox(sld, name, Inches(0.4), top + Inches(0.05), Inches(6.5), Inches(0.27), fs=11, color=C_BLACK)
        rect(sld, Inches(7.85), top, Inches(1.8), gr_h - Inches(0.03), C_TEAL)
        txbox(sld, mode, Inches(7.85), top + Inches(0.05), Inches(1.8), Inches(0.27),
              fs=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    slide_num_label(sld, num)


def slide_12_setup(prs, num):
    """セットアップ"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "セットアップ")
    bottom_line(sld)

    y = content_start_y()

    # conda環境作成コマンド
    txbox(sld, "■ conda 環境作成",
          Inches(0.3), y, W - Inches(0.6), Inches(0.32), fs=14, bold=True, color=C_HEADER)
    cmds = [
        "conda create -n met_env_310 python=3.10",
        "conda activate met_env_310",
        "conda install -c conda-forge pygrib xarray metpy matplotlib cartopy requests",
        "pip install beautifulsoup4 python-pptx",
    ]
    cmd_y = y + Inches(0.36)
    rect(sld, Inches(0.3), cmd_y, W - Inches(0.6), Inches(len(cmds) * 0.42 + 0.1),
         RGBColor(0xF0, 0xF0, 0xF8))
    for i, cmd in enumerate(cmds):
        txbox(sld, cmd, Inches(0.45), cmd_y + Inches(0.05) + i * Inches(0.42),
              W - Inches(0.9), Inches(0.38), fs=12, color=C_ACCENT2, italic=True)

    # ライブラリ一覧
    lib_y = cmd_y + Inches(len(cmds) * 0.42 + 0.25)
    txbox(sld, "■ 主要ライブラリ",
          Inches(0.3), lib_y, W - Inches(0.6), Inches(0.32), fs=14, bold=True, color=C_HEADER)

    libs = [
        ("pygrib",        "GRIB2ファイル読み込み・変数抽出"),
        ("xarray",        "多次元データセット管理"),
        ("metpy",         "渦度・発散・相当温位等の気象計算"),
        ("matplotlib",    "天気図描画"),
        ("cartopy",       "地図投影・海岸線描画"),
        ("requests",      "HTTPダウンロード（RISHサーバー・ECMWF）"),
        ("beautifulsoup4","HTMLパース（ファイルリスト取得）"),
        ("python-pptx",   "PowerPointファイル自動生成"),
    ]
    row_h = Inches(0.42)
    for i, (lib, desc) in enumerate(libs):
        top = lib_y + Inches(0.36) + i * row_h
        bg = C_ACCENT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), top, W - Inches(0.6), row_h - Inches(0.04), bg)
        txbox(sld, lib, Inches(0.35), top + Inches(0.07), Inches(1.8), Inches(0.32),
              fs=11, bold=True, color=C_HEADER, italic=True)
        txbox(sld, desc, Inches(2.2), top + Inches(0.07), Inches(7.4), Inches(0.32),
              fs=11, color=C_BLACK)

    slide_num_label(sld, num)


def slide_13_output(prs, num):
    """出力ファイルとサンプル"""
    sld = new_slide(prs)
    set_bg(sld, C_BG)
    header_bar(sld, "出力ファイルとサンプル")
    bottom_line(sld)

    y = content_start_y()

    # 出力命名規則
    txbox(sld, "■ 出力ファイル命名規則: output/{YYYYMMDDHH}_FT{FFF}h_{種別}.png",
          Inches(0.3), y, W - Inches(0.6), Inches(0.32), fs=13, bold=True, color=C_HEADER)

    examples = [
        ("GSM 500hPa渦度",  "2026041200_FT000h_500hPa_Height_VORT.png"),
        ("GSM FAX57",        "2026041200_FT000h_GSM_Fax57.png"),
        ("GSM FAX78",        "2026041200_FT000h_GSM_Fax78.png"),
        ("GSM 地上気圧",     "2026041200_FT000h_GSM_SurfacePressure.png"),
        ("GSM 850hPa EPT",   "2026041200_FT000h_GSM_850hPa_EPT.png"),
        ("ECM 500hPa渦度",   "2026041200_FT000h_ECM_500hPa_Height_VORT.png"),
        ("ECM 地上気圧",     "2026041200_FT000h_ECM_SurfacePressure.png"),
    ]
    ex_y = y + Inches(0.36)
    row_h = Inches(0.42)
    for i, (label, fname) in enumerate(examples):
        top = ex_y + i * row_h
        bg = C_ACCENT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), top, W - Inches(0.6), row_h - Inches(0.04), bg)
        txbox(sld, label, Inches(0.35), top + Inches(0.07), Inches(1.8), Inches(0.32),
              fs=11, bold=True, color=C_HEADER)
        txbox(sld, fname, Inches(2.2), top + Inches(0.07), Inches(7.4), Inches(0.32),
              fs=11, color=C_ACCENT2, italic=True)

    # サンプル
    samp_y = ex_y + len(examples) * row_h + Inches(0.2)
    txbox(sld, "■ samples/ ディレクトリ",
          Inches(0.3), samp_y, W - Inches(0.6), Inches(0.32), fs=13, bold=True, color=C_HEADER)
    txbox(sld, "全出力種別のサンプルPNG 14枚 + PowerPoint 1ファイルを収録（2026-04-12 12UTC 初期値）\n"
               "sample_tenkizu_2026041212.pptx: 17スライドのPowerPointサンプルも収録",
          Inches(0.3), samp_y + Inches(0.36), W - Inches(0.6), Inches(0.6),
          fs=12, color=C_BLACK)

    slide_num_label(sld, num)


def slide_14_summary(prs, num):
    """まとめ"""
    sld = new_slide(prs)
    set_bg(sld, C_HEADER)  # 濃紺背景でまとめを強調

    # タイトル
    txbox(sld, "まとめ",
          Inches(0.5), Inches(0.2), W - Inches(1), Inches(0.75),
          fs=32, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # 区切り線
    rect(sld, Inches(0.5), Inches(0.95), W - Inches(1), Inches(0.05), C_WHITE)

    points = [
        ("GSM・ECMWF 対応",
         "気象庁GSM（RISH経由・過去データも無償）と\nECMWF（Open Data・最新5日分無償）の両方に対応"),
        ("描画スクリプト 16本",
         "500hPa天気図・FAX57/78・EPT・ジェット・鉛直断面など\n高層から地上まで網羅的に対応"),
        ("自動化",
         "run_gsm_auto.py / run_ecm_auto.py で最新データを自動検索し\n全スクリプトを一括実行"),
        ("レポート・PowerPoint 自動生成",
         "Markdownレポート（GitHub push対応）と\nPowerPointファイルの自動生成機能を装備"),
        ("拡張性",
         "引数仕様を全スクリプトで統一。任意気圧面・複数FT・\nECMオプション(tcwv/tp)など柔軟に対応"),
    ]

    y = Inches(1.15)
    row_h = Inches(1.05)
    for i, (title, desc) in enumerate(points):
        top = y + i * row_h
        rect(sld, Inches(0.3), top, Inches(0.08), row_h - Inches(0.1), C_TEAL)
        txbox(sld, title, Inches(0.5), top + Inches(0.03), W - Inches(0.7), Inches(0.35),
              fs=15, bold=True, color=C_ACCENT)
        txbox(sld, desc, Inches(0.5), top + Inches(0.38), W - Inches(0.7), Inches(0.55),
              fs=13, color=C_WHITE)

    slide_num_label(sld, num)


# ============================================================
# メイン
# ============================================================

def main():
    output = Path("/Users/masahiro/projects/ppt_auto/tenkizu_overview.pptx")

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_title(prs)
    slide_02_overview(prs, 2)
    slide_03_gsm(prs, 3)
    slide_04_ecmwf(prs, 4)
    slide_05_comparison(prs, 5)
    slide_06_scripts_gsm(prs, 6)
    slide_07_scripts_ecm(prs, 7)
    slide_08_args_and_bulk(prs, 8)
    slide_09_auto(prs, 9)
    slide_10_reports(prs, 10)
    slide_11_pptx_gen(prs, 11)
    slide_12_setup(prs, 12)
    slide_13_output(prs, 13)
    slide_14_summary(prs, 14)

    prs.save(output)
    print(f"生成完了: {output}")
    print(f"総スライド数: {len(prs.slides)}枚")


if __name__ == "__main__":
    main()
