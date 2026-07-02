"""
PowerPoint定型スライド自動生成スクリプト
使い方: python3 make_pptx.py config.json  （またはMarkdownファイルを直接指定: python3 make_pptx.py slides.md）
"""

import json
import math
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR

EMU_PER_PT = 12700

# コンテンツ領域（タイトルバー直下から下部バー直上まで）
CONTENT_LEFT   = Inches(0.8)
CONTENT_WIDTH  = Inches(11.7)
CONTENT_TOP    = Inches(1.5)


# スライドサイズ（16:9）
SLIDE_WIDTH  = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# カラー定義
COLOR_ACCENT  = RGBColor(0x1F, 0x49, 0x7D)  # 濃紺（タイトルバー）
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK    = RGBColor(0x26, 0x26, 0x26)   # 本文テキスト
COLOR_LIGHT   = RGBColor(0xF2, 0xF2, 0xF2)  # 背景薄グレー
COLOR_CODE_BG = RGBColor(0xE8, 0xE8, 0xE8)  # コードブロック背景（本文背景より少し濃いグレー）


def set_bg(slide, color):
    """スライド背景色を設定する"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text_width_pt(text, font_size_pt):
    """テキストの概算幅をポイント単位で計算する
    （全角文字はフォントサイズと同程度、半角英数字は約0.55倍として近似する）"""
    width = 0.0
    for ch in text:
        width += font_size_pt * (0.55 if ord(ch) < 128 else 1.0)
    return width


def _estimate_wrapped_lines(text, font_size_pt, width_emu):
    """word_wrap時にテキストが折り返されて何行になるか概算する"""
    if not text:
        return 1
    width_pt = width_emu / EMU_PER_PT
    if width_pt <= 0:
        return 1
    return max(1, math.ceil(_text_width_pt(text, font_size_pt) / width_pt))


def _estimate_text_block_height(text, font_size_pt, width_emu, line_spacing=1.25, pad=Inches(0.15)):
    """\\n区切りの複数行テキストがword_wrap時に必要とする概算の高さ(EMU)を返す"""
    lines = text.split("\n")
    total_lines = sum(_estimate_wrapped_lines(line, font_size_pt, width_emu) for line in lines)
    line_height_emu = font_size_pt * line_spacing * EMU_PER_PT
    return Emu(int(total_lines * line_height_emu)) + pad


def _mono_text_width_pt(text, font_size_pt, char_width_ratio=0.62):
    """等幅フォントのテキスト幅を概算する
    （半角文字はfont_size*char_width_ratio、全角文字はfont_size分として近似する）"""
    width = 0.0
    for ch in text:
        width += font_size_pt * (char_width_ratio if ord(ch) < 128 else 1.0)
    return width


def _estimate_mono_wrapped_lines(text, font_size_pt, width_emu, char_width_ratio=0.62):
    """等幅フォントのテキストがword_wrap時に何行になるか概算する"""
    if not text:
        return 1
    width_pt = width_emu / EMU_PER_PT
    if width_pt <= 0:
        return 1
    return max(1, math.ceil(_mono_text_width_pt(text, font_size_pt, char_width_ratio) / width_pt))


def _estimate_code_block_height(text, font_size_pt, width_emu, line_spacing=1.15, pad=Inches(0.2)):
    """コードブロック（等幅フォント）が必要とする概算の高さ(EMU)を返す"""
    lines = text.split("\n")
    total_lines = sum(_estimate_mono_wrapped_lines(line, font_size_pt, width_emu) for line in lines)
    line_height_emu = font_size_pt * line_spacing * EMU_PER_PT
    return Emu(int(total_lines * line_height_emu)) + pad


def add_rect(slide, left, top, width, height, color):
    """塗りつぶし矩形を追加する"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # 枠線なし
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=COLOR_DARK,
                 align=PP_ALIGN.LEFT, wrap=True, autofit=False, font_name="メイリオ"):
    """テキストボックスを追加する（\\n はPowerPointの段落として分割する。
    run.textに\\nをそのまま渡すとXML上は改行文字が入るだけでPowerPointは
    改行として描画しないため、行ごとに別段落を作る必要がある）。
    autofit=Trueの場合、ボックスに収まらない量のテキストが入ったときPowerPoint側で
    自動的にフォントサイズを縮小する（python-pptxはフォントメトリクスを持たないため
    事前計算はできず、PowerPointの「自動調整」機能に委ねる）"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if autofit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox


def make_title_slide(prs, cfg):
    """表紙スライドを作成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
    set_bg(slide, COLOR_LIGHT)

    # 上部アクセントバー
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.15), COLOR_ACCENT)

    # 中央タイトル背景（タイトルが\nで複数行になる場合に備えて高さを確保）
    add_rect(slide, Inches(1), Inches(1.8), Inches(11.33), Inches(2.7), COLOR_ACCENT)

    # タイトル
    add_text_box(
        slide, cfg["title"],
        Inches(1.3), Inches(2.0), Inches(10.8), Inches(1.6),
        font_size=40, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER
    )

    # サブタイトル
    add_text_box(
        slide, cfg.get("subtitle", ""),
        Inches(1.3), Inches(3.6), Inches(10.8), Inches(0.8),
        font_size=24, color=COLOR_WHITE, align=PP_ALIGN.CENTER
    )

    # 作成者・日付
    info = f"{cfg.get('date', '')}　{cfg.get('author', '')}"
    add_text_box(
        slide, info,
        Inches(1), Inches(5.5), Inches(11.33), Inches(0.5),
        font_size=16, color=COLOR_DARK, align=PP_ALIGN.RIGHT
    )

    # 下部アクセントバー
    add_rect(slide, 0, Inches(7.2), SLIDE_WIDTH, Inches(0.3), COLOR_ACCENT)


def make_toc_slide(prs, cfg):
    """目次スライドを作成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_LIGHT)

    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(1.2), COLOR_ACCENT)
    add_text_box(
        slide, "目次",
        Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
        font_size=32, bold=True, color=COLOR_WHITE
    )

    contents = [s["title"] for s in cfg.get("slides", [])]
    body = "\n".join(f"  {i+1}.  {t}" for i, t in enumerate(contents))
    add_text_box(
        slide, body,
        Inches(1.0), Inches(1.5), Inches(11), Inches(5.5),
        font_size=22, color=COLOR_DARK, autofit=True
    )

    add_rect(slide, 0, Inches(7.2), SLIDE_WIDTH, Inches(0.3), COLOR_ACCENT)


def add_quote_box(slide, text, top):
    """引用（blockquote）を左アクセント線付きの帯として追加する"""
    height = Inches(0.9)
    add_rect(slide, CONTENT_LEFT, top, CONTENT_WIDTH, height, COLOR_LIGHT)
    add_rect(slide, CONTENT_LEFT, top, Inches(0.06), height, COLOR_ACCENT)
    add_text_box(
        slide, text,
        CONTENT_LEFT + Inches(0.25), top + Inches(0.08), CONTENT_WIDTH - Inches(0.4), height - Inches(0.16),
        font_size=16, color=COLOR_ACCENT, autofit=True
    )
    return top + height + Inches(0.15)


def add_images_row(slide, image_paths, top, max_height=Inches(4.2)):
    """画像を横並びで挿入する（アスペクト比は維持したまま、全画像の高さを揃える）"""
    n = len(image_paths)
    gap = Inches(0.2)
    img_width = Emu(int((CONTENT_WIDTH - gap * (n - 1)) / n))

    # 各画像のアスペクト比を先に調べ、列幅いっぱいに収めたときの高さの最小値を
    # 全画像共通の高さとする（アスペクト比がバラバラだと高さが不揃いになるため）
    aspects = []
    for path in image_paths:
        pic = slide.shapes.add_picture(path, 0, 0, width=img_width)
        aspects.append(pic.width / pic.height)
        pic._element.getparent().remove(pic._element)

    common_height = min(Emu(int(img_width / a)) for a in aspects)
    common_height = min(common_height, max_height)

    left = CONTENT_LEFT
    for path, aspect in zip(image_paths, aspects):
        width = Emu(int(common_height * aspect))
        pic_left = int(left + (img_width - width) / 2)
        slide.shapes.add_picture(path, pic_left, top, width=width, height=common_height)
        left += img_width + gap
    return top + common_height + Inches(0.15)


def add_table_shape(slide, table_data, top, max_height=Inches(4.5), font_size=14):
    """Markdownテーブルを本物のPowerPoint表として挿入する
    （箇条書きへの機械変換は行単位の意味関係が失われて読みにくくなるため、
    構造をそのまま保てるネイティブの表オブジェクトを使う）。
    各行の高さは固定値ではなく、セルの実際の文字量（折り返し行数）から概算する
    （固定値だと長いセル内容がPowerPoint側で自動的に行が拡張され、想定より
    はみ出してしまうため）"""
    header = table_data["header"]
    rows_data = table_data["rows"]
    n_cols = len(header)
    n_rows = 1 + len(rows_data)

    col_width = Emu(int(CONTENT_WIDTH / n_cols))
    cell_text_width = max(Emu(1), col_width - Inches(0.16))  # 左右マージン分を引く
    min_row_height = Inches(0.35)

    def _row_height(cells):
        max_lines = max(
            (_estimate_wrapped_lines(c, font_size, cell_text_width) for c in cells),
            default=1,
        )
        line_height = Emu(int(font_size * 1.3 * EMU_PER_PT))
        return max(min_row_height, Emu(max_lines * line_height) + Inches(0.06))

    row_heights = [_row_height(header)] + [_row_height(row) for row in rows_data]
    table_height = Emu(sum(row_heights))

    # 概算の合計高さが確保できる範囲を超える場合は、行の高さを一律縮小して収める
    # （それでも内容量が多すぎる場合はページ溢れとなるため、手動での内容整理が必要）
    if table_height > max_height:
        scale = max_height / table_height
        row_heights = [Emu(max(int(min_row_height * 0.5), int(h * scale))) for h in row_heights]
        table_height = Emu(sum(row_heights))

    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, CONTENT_LEFT, top, CONTENT_WIDTH, table_height
    )
    table = graphic_frame.table
    table.first_row = True
    table.horz_banding = False

    for c in range(n_cols):
        table.columns[c].width = col_width
    for r, h in enumerate(row_heights):
        table.rows[r].height = h

    def _fill_cell(cell, text, bold, font_color, bg_color):
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = font_color
                run.font.name = "メイリオ"

    for c, text in enumerate(header):
        _fill_cell(table.cell(0, c), text, True, COLOR_WHITE, COLOR_ACCENT)

    for r, row in enumerate(rows_data, start=1):
        for c in range(n_cols):
            text = row[c] if c < len(row) else ""
            _fill_cell(table.cell(r, c), text, False, COLOR_DARK, COLOR_WHITE)

    return top + table_height + Inches(0.2)


def add_code_box(slide, text, top, max_height=Inches(3.0), font_size=13):
    """コードブロック（```...```で囲まれた部分）を等幅フォントのテキストボックスとして表示する。
    ディレクトリツリーなどのファイル構成図はインデント・記号の位置関係が重要なため、
    崩れないよう等幅フォント（Courier New）でそのまま表示する方針にした
    （以前は「スライドに不向き」として省略していたが、構成図等の情報が失われるとの
    指摘を受けて、省略ではなく表示する方式に変更した）。
    高さは実際の行数・折り返しから概算し、収まらない場合はテーブルと同様に
    一律縮小して他の要素との重なりを防ぐ"""
    inner_width = CONTENT_WIDTH - Inches(0.5)
    height = _estimate_code_block_height(text, font_size, inner_width)

    if height > max_height:
        # 概算の高さが確保できる範囲を超える場合は、文字が読める最小サイズまで
        # フォントを段階的に縮小して収める（テーブルの行縮小と同じ考え方）
        for smaller_size in (font_size - 1, font_size - 2, font_size - 3):
            if smaller_size < 8:
                break
            candidate = _estimate_code_block_height(text, smaller_size, inner_width)
            if candidate <= max_height:
                font_size = smaller_size
                height = candidate
                break
        height = min(height, max_height)

    add_rect(slide, CONTENT_LEFT, top, CONTENT_WIDTH, height, COLOR_CODE_BG)
    add_rect(slide, CONTENT_LEFT, top, Inches(0.06), height, COLOR_ACCENT)
    add_text_box(
        slide, text,
        CONTENT_LEFT + Inches(0.3), top + Inches(0.1), inner_width - Inches(0.2), height - Inches(0.2),
        font_size=font_size, color=COLOR_DARK, font_name="Courier New"
    )
    return top + height + Inches(0.2)


def make_content_slide(prs, slide_cfg):
    """本文スライドを作成する（本文・画像・引用・表・コードブロックの組み合わせに対応）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_LIGHT)

    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(1.2), COLOR_ACCENT)
    add_text_box(
        slide, slide_cfg["title"],
        Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
        font_size=32, bold=True, color=COLOR_WHITE
    )

    body = slide_cfg.get("body", "")
    images = slide_cfg.get("images") or ([slide_cfg["image"]] if slide_cfg.get("image") else [])
    tables = slide_cfg.get("tables") or []
    code_blocks = slide_cfg.get("code_blocks") or []
    quote = slide_cfg.get("quote", "")

    cursor = CONTENT_TOP

    if body:
        # 画像・表・コードブロックがある場合はフォントをやや小さくする。高さは固定値ではなく
        # 実際の文字量から概算する（固定値だと本文が長い場合に下の要素と重なってしまうため）
        font_size = 20 if (images or tables or code_blocks) else 22
        body_height = max(Inches(0.5), _estimate_text_block_height(body, font_size, CONTENT_WIDTH))
        body_height = min(body_height, Inches(5.5))  # 他の要素のための余白は確保する
        add_text_box(
            slide, body,
            CONTENT_LEFT, cursor, CONTENT_WIDTH, body_height,
            font_size=font_size, color=COLOR_DARK, autofit=True
        )
        cursor += body_height

    for table_data in tables:
        remaining = max(Inches(0.3), Inches(7.0) - cursor - (Inches(1.05) if quote else 0))
        cursor = add_table_shape(slide, table_data, cursor, max_height=remaining)

    for code_text in code_blocks:
        remaining = max(Inches(0.3), Inches(7.0) - cursor - (Inches(1.05) if quote else 0))
        cursor = add_code_box(slide, code_text, cursor, max_height=remaining)

    if images:
        remaining = max(Inches(0.3), Inches(7.0) - cursor - (Inches(1.05) if quote else 0))
        cursor = add_images_row(slide, images, cursor, max_height=remaining)

    if quote:
        cursor = min(cursor, Inches(6.0))
        add_quote_box(slide, quote, cursor)

    add_rect(slide, 0, Inches(7.2), SLIDE_WIDTH, Inches(0.3), COLOR_ACCENT)


def make_summary_slide(prs, title, body):
    """まとめスライドを作成する（背景色を変えて強調）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_ACCENT)

    add_text_box(
        slide, title,
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.9),
        font_size=36, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER
    )

    add_rect(slide, Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.05), COLOR_WHITE)

    add_text_box(
        slide, body,
        Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5),
        font_size=22, color=COLOR_WHITE, autofit=True
    )


def _strip_inline_md(text: str) -> str:
    """行内の**太字**・*イタリック*・`コード`記法の記号だけを取り除く"""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text


def _split_table_cells(line: str) -> list:
    """Markdownテーブルの1行をセルのリストに分解する"""
    inner = line.strip()
    if inner.startswith("|"):
        inner = inner[1:]
    if inner.endswith("|"):
        inner = inner[:-1]
    return [_strip_inline_md(c.strip()) for c in inner.split("|")]


def _is_table_separator(cells: list) -> bool:
    """Markdownテーブルの区切り行（|---|---|等）かどうか判定する"""
    return bool(cells) and all(re.fullmatch(r":?-+:?", c) for c in cells)


def parse_markdown(md_path: Path) -> dict:
    """Markdownファイルをconfig.jsonと同じ構造のdictに変換する（config.json不要でMarkdownから直接生成するための入口）。

    対応フォーマット（Marpに似せたシンプルな記法）:
      ---
      title: タイトル
      subtitle: サブタイトル
      author: 作成者
      date: 日付
      ---

      ## スライドタイトル

      本文テキスト
      - 箇条書き（・付きに自動変換される）

      > 引用・強調したい一言（quoteフィールドになる）

      ![説明](image.png)（imagesフィールドになる。複数書けば横並び）

      | 列1 | 列2 |
      |---|---|
      | a | b |
      （本物のPowerPoint表に変換される）

      ```
      project/
      ├── src/
      └── README.md
      ```
      （コードブロックは等幅フォントのテキストボックスとして表示される）

      ## まとめ <!-- summary -->

      まとめ本文

    フロントマター（---...---）は丸ごと省略可能。titleが指定されていない場合、
    本文の最初の"# "見出し（H1）、なければ最初の"## "見出し（H2）を自動でタイトルに
    使う。subtitle/author/dateは元々省略可（空欄）。これにより、README.md等の
    既存Markdownをそのまま渡しても動作する。
    """
    content = md_path.read_text(encoding="utf-8")
    cfg = {"title": "", "subtitle": "", "author": "", "date": "", "slides": []}

    # フロントマター（---...---）を解析（フロントマター自体が省略されていてもよい）
    fm_match = re.match(r"^---\s*\n([\s\S]*?)\n---\s*\n?", content)
    if fm_match:
        for line in fm_match.group(1).split("\n"):
            if ":" in line:
                k, _, v = line.partition(":")
                k, v = k.strip(), v.strip()
                if k in cfg:
                    cfg[k] = v
        content = content[fm_match.end():]

    # コードブロック（```...```）内の行は見出し検出の対象から除外する
    # （bashコメント"# ..."やPythonコメント"# ..."を見出しと誤認識しないようにするため）
    lines = content.split("\n")
    fenced = set()
    in_fence = False
    for idx, ln in enumerate(lines):
        if ln.strip().startswith("```"):
            in_fence = not in_fence
            fenced.add(idx)
            continue
        if in_fence:
            fenced.add(idx)

    def _first_heading(level_pattern):
        for idx, ln in enumerate(lines):
            if idx in fenced:
                continue
            m = re.match(level_pattern, ln)
            if m:
                return m.group(1).strip()
        return None

    # titleが未指定の場合、本文の見出しから自動取得する
    # （最初の"# "見出し、なければ最初の"## "見出しをタイトルとして使う。
    # subtitle/author/dateは元々省略可＝空欄でよいため、これでフロントマター自体が不要になる）
    if not cfg["title"]:
        cfg["title"] = _first_heading(r"^#\s+(.+)$") or _first_heading(r"^##\s+(.+)$") or ""

    # H1/H2見出しでスライドに分割する（コードブロック内の"#"はここでは見出しとみなさない。
    # H3以下はスライド内の小見出しとして扱う）
    heading_positions = []
    for idx, ln in enumerate(lines):
        if idx in fenced:
            continue
        m = re.match(r"^#{1,2}\s+(.+)$", ln)
        if m:
            heading_positions.append((idx, m.group(1)))

    slide_sections = []
    for k, (idx, raw_title) in enumerate(heading_positions):
        body_start = idx + 1
        body_end = heading_positions[k + 1][0] if k + 1 < len(heading_positions) else len(lines)
        slide_sections.append((raw_title, "\n".join(lines[body_start:body_end])))

    for raw_title, raw_body in slide_sections:
        raw_title = raw_title.strip()
        raw_body = raw_body.strip()

        slide_type = "content"
        if "<!-- summary -->" in raw_title:
            slide_type = "summary"
        raw_title = re.sub(r"<!--.*?-->", "", raw_title).strip()

        body_lines = []
        quote_lines = []
        image_paths = []
        tables = []
        code_blocks = []
        in_code_block = False
        code_block_lines = None
        table_header = None
        table_rows = None
        table_pending_header = False

        def _flush_table():
            """表の読み取り中だった場合、確定した表をtablesに積む"""
            if table_header is not None and table_rows:
                tables.append({"header": table_header, "rows": table_rows})

        for line in raw_body.split("\n"):
            stripped = line.strip()

            # コードブロック（```...```）は等幅フォントの表示用に本文とは別に集める
            # （インデント等の位置関係を保つため、strip前の行をそのまま保持する）
            if stripped.startswith("```"):
                if in_code_block:
                    if code_block_lines:
                        code_blocks.append("\n".join(code_block_lines).rstrip("\n"))
                    code_block_lines = None
                    in_code_block = False
                else:
                    in_code_block = True
                    code_block_lines = []
                continue
            if in_code_block:
                code_block_lines.append(line.rstrip())
                continue

            # Markdownテーブル（| a | b |）は本物のPowerPoint表として扱う
            if stripped.startswith("|") and stripped.endswith("|") and len(stripped) >= 2:
                cells = _split_table_cells(stripped)
                if table_pending_header:
                    if _is_table_separator(cells):
                        table_pending_header = False
                        table_rows = []
                        continue
                    # 直前の行はヘッダーではなく単なるデータ行だった
                    _flush_table()
                    table_header = cells
                    table_rows = None
                    table_pending_header = True
                    continue
                elif table_rows is not None:
                    table_rows.append(cells)
                    continue
                else:
                    table_header = cells
                    table_pending_header = True
                    continue
            else:
                if table_pending_header:
                    # ヘッダー候補だけで表として確定しなかった場合は普通の行として扱う
                    body_lines.append("・" + " / ".join(table_header))
                else:
                    _flush_table()
                table_header = None
                table_rows = None
                table_pending_header = False

            # 水平線（---, ___, ***）は無視する
            if re.fullmatch(r"[-_*]{3,}", stripped):
                continue

            m_img = re.match(r"!\[.*?\]\((.+?)\)", stripped)
            if m_img:
                image_paths.append(m_img.group(1))
                continue

            if stripped.startswith(">"):
                quote_lines.append(stripped.lstrip(">").strip())
                continue

            # H3〜H6はスライドを分けず、記号を外して小見出し行として扱う
            m_subhead = re.match(r"^#{3,6}\s+(.*)", stripped)
            if m_subhead:
                body_lines.append(_strip_inline_md(m_subhead.group(1)))
                continue

            m_bullet = re.match(r"^[-*+]\s+(.*)", stripped)
            prefix, text = ("・", m_bullet.group(1)) if m_bullet else ("", stripped)
            body_lines.append(prefix + _strip_inline_md(text))

        _flush_table()
        # 閉じフェンスが無いまま本文が終わった場合も、それまでの内容を確定させる
        if in_code_block and code_block_lines:
            code_blocks.append("\n".join(code_block_lines).rstrip("\n"))

        slide_cfg = {
            "type": slide_type,
            "title": raw_title,
            "body": "\n".join(body_lines).strip("\n"),
        }
        if image_paths:
            # 画像パスはMarkdownファイルからの相対パスとして解決する
            slide_cfg["images"] = [
                str((md_path.parent / p).resolve()) if not Path(p).is_absolute() else p
                for p in image_paths
            ]
        if quote_lines:
            slide_cfg["quote"] = _strip_inline_md(" ".join(quote_lines))
        if tables:
            slide_cfg["tables"] = tables
        if code_blocks:
            slide_cfg["code_blocks"] = code_blocks

        cfg["slides"].append(slide_cfg)

    return cfg


MARKDOWN_TEMPLATE_HELP = """\
使い方: python3 make_pptx.py config.json
       python3 make_pptx.py slides.md

Markdownファイルの記法:

---
title: プレゼンタイトル
subtitle: サブタイトル
author: 作成者名
date: 2026年7月
---

## スライドタイトル

本文テキスト
- 箇条書き1
- 箇条書き2

> 引用・強調したい一言

![説明](画像のパス)

| 列1 | 列2 |
|---|---|
| a | b |

```
project/
├── src/
└── README.md
```

## まとめ <!-- summary -->

まとめ本文

補足: フロントマター（---...---）は省略可能。titleを省略した場合、本文の
最初の "# " 見出し（なければ最初の "## " 見出し）を自動でタイトルに使う。
subtitle/author/dateも省略時は空欄になるだけなので、既存のMarkdownファイル
（README.md等）をそのまま渡しても動作する。
表（| ... |）は本物のPowerPoint表に、コードブロック（```...```）は等幅
フォントのテキストボックス（ディレクトリツリー等のファイル構成表示向け）に
それぞれ変換される。いずれも他の要素との重なりを防ぐため、内容量に応じて
高さ・フォントサイズを自動調整するが、内容が多すぎる場合は手動での調整が必要。
"""


def main():
    if len(sys.argv) < 2:
        print("使い方: python3 make_pptx.py config.json")
        print("       python3 make_pptx.py slides.md")
        sys.exit(1)

    if sys.argv[1] == "?":
        print(MARKDOWN_TEMPLATE_HELP)
        sys.exit(0)

    config_path = Path(sys.argv[1])
    if not config_path.exists():
        print(f"エラー: {config_path} が見つかりません")
        sys.exit(1)

    if config_path.suffix.lower() == ".md":
        cfg = parse_markdown(config_path)
    else:
        with open(config_path, encoding="utf-8") as f:
            cfg = json.load(f)

    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 表紙
    make_title_slide(prs, cfg)

    # 目次（スライドが2枚以上ある場合のみ）
    if len(cfg.get("slides", [])) >= 2:
        make_toc_slide(prs, cfg)

    # 本文・まとめ
    for slide_cfg in cfg.get("slides", []):
        if slide_cfg.get("type") == "summary":
            make_summary_slide(prs, slide_cfg["title"], slide_cfg["body"])
        else:
            make_content_slide(prs, slide_cfg)

    out_path = config_path.with_name(config_path.stem + ".pptx")
    prs.save(out_path)
    print(f"生成完了: {out_path}")


if __name__ == "__main__":
    main()
